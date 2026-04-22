from datetime import datetime, timedelta
from pathlib import Path
import tempfile
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

import pandas as pd
from pptx import Presentation
from lxml import etree

from logic.post_mortem import PostMortemAnalyzer
from logic.powerpoint_report import (
    MillRelinePowerPointReport,
    apply_placeholder_replacements,
    build_default_replacements,
)
from logic.report_export import _build_delay_summary, _format_delay_summary_for_table, build_excel_report_buffer
from models.delay import Delay, DelayType
from models.phase import Phase
from models.project import Project
from models.task import Task


def _build_project() -> Project:
    phase_1 = Phase(name="Tear Out")
    phase_1.add_task(
        Task(
            name="Task 1",
            start_date=datetime(2026, 2, 1, 7, 0),
            end_date=datetime(2026, 2, 1, 9, 0),
            actual_start=datetime(2026, 2, 1, 7, 0),
            actual_end=datetime(2026, 2, 1, 10, 0),
            note="Load level",
        )
    )
    phase_1.add_task(
        Task(
            name="Task 2",
            start_date=datetime(2026, 2, 1, 9, 0),
            end_date=datetime(2026, 2, 1, 12, 0),
            actual_start=datetime(2026, 2, 1, 10, 0),
            actual_end=datetime(2026, 2, 1, 13, 0),
            note="",
        )
    )

    phase_2 = Phase(name="Install")
    phase_2.add_task(
        Task(
            name="Task 3",
            start_date=datetime(2026, 2, 2, 7, 0),
            end_date=datetime(2026, 2, 2, 9, 0),
            actual_start=datetime(2026, 2, 2, 7, 0),
            actual_end=datetime(2026, 2, 2, 8, 30),
            note="Finished early",
        )
    )

    project = Project(name="HVC SAG Mill")
    project.add_phase(phase_1)
    project.add_phase(phase_2)
    return project


def test_apply_placeholder_replacements_supports_split_phase_tokens():
    text = "Phase < | phase_num | >: < | phase_name | > - Timeline"
    replaced = apply_placeholder_replacements(
        text,
        {
            "phase_num": 3,
            "phase_name": "Install",
        },
    )

    assert replaced == "Phase 3: Install - Timeline"


def test_apply_placeholder_replacements_strips_powerpoint_control_escape_sequences():
    text = "Phase < | phase_num | >: < | phase_name | > - Delayed Tasks"
    replaced = apply_placeholder_replacements(
        text,
        {
            "phase_num": 2,
            "phase_name": "_x000B_Stripping FH & Shell_x000B_",
        },
    )

    assert replaced == "Phase 2: Stripping FH & Shell - Delayed Tasks"


def test_replace_phase_title_paragraph_preserves_delayed_task_title_structure():
    prs = Presentation(str(Path(__file__).resolve().parents[1] / "src" / "assets" / "mill_reline_template.pptx"))
    paragraph = prs.slides[15].shapes[5].text_frame.paragraphs[0]

    replaced = MillRelinePowerPointReport._replace_phase_title_paragraph(
        paragraph,
        {
            "phase_num": 2,
            "phase_name": "Stripping FH & Shell",
        },
    )

    runs = [run.text for run in paragraph.runs]
    assert replaced is True
    assert runs[0] == "Phase "
    assert runs[1] == "2"
    assert runs[2] == ": "
    assert runs[3] == "Stripping FH & Shell"
    assert runs[4] == " - "
    assert runs[6] == "Delayed Tasks"


def test_replace_phase_title_paragraph_noops_without_phase_values():
    prs = Presentation(str(Path(__file__).resolve().parents[1] / "src" / "assets" / "mill_reline_template.pptx"))
    paragraph = prs.slides[15].shapes[5].text_frame.paragraphs[0]
    original_runs = [run.text for run in paragraph.runs]

    replaced = MillRelinePowerPointReport._replace_phase_title_paragraph(
        paragraph,
        {
            "PROJECT_TITLE": "Test Project",
        },
    )

    assert replaced is False
    assert [run.text for run in paragraph.runs] == original_runs


def test_build_default_replacements_uses_project_dates():
    project = _build_project()

    replacements = build_default_replacements(
        project,
        created_at=datetime(2026, 3, 13, 9, 0),
    )

    assert replacements["PROJECT_TITLE"] == "HVC SAG Mill"
    assert replacements["MONTH"] == "February"
    assert replacements["YEAR"] == "2026"
    assert replacements["START_DAY"] == "01"
    assert replacements["END_DAY"] == "02"
    assert replacements["TOTAL_DELAY_HOURS"] == -0.5


def test_phase_delay_summary_builds_reporting_table():
    project = _build_project()

    summary = PostMortemAnalyzer.phase_delay_summary(project)

    assert list(summary.columns) == [
        "No.",
        "PHASE",
        "PHASE DELAY (HRS)",
        "CUMULATIVE DELAY (HRS)",
        "RATIONALE",
    ]
    assert summary["PHASE"].tolist() == ["Tear Out", "Install"]
    assert summary["PHASE DELAY (HRS)"].tolist() == [1.0, -0.5]
    assert summary["CUMULATIVE DELAY (HRS)"].tolist() == [1.0, -0.5]


def test_build_phase_sections_matches_project_shape():
    project = _build_project()

    sections = MillRelinePowerPointReport.build_phase_sections(
        project,
        phase_delay_rows=1,
    )

    assert len(sections) == 2
    assert sections[0].replacements["phase_num"] == 1
    assert sections[0].replacements["phase_name"] == "Tear Out"
    assert sections[1].replacements["phase_num"] == 2
    assert len(sections[0].delayed_tasks) == 1


def test_build_delay_summary_groups_delay_types():
    delays = [
        Delay(
            id="1",
            project_id="p1",
            delay_type=DelayType.PREPARATION,
            duration_minutes=60,
            description="Prep",
            start_dt=datetime(2026, 2, 1, 7, 0),
            end_dt=datetime(2026, 2, 1, 8, 0),
            shift_assignment_id=None,
            created_by="me",
            created_at=datetime(2026, 2, 1, 8, 0),
            updated_at=None,
            updated_by=None,
        ),
        Delay(
            id="2",
            project_id="p1",
            delay_type=DelayType.PREPARATION,
            duration_minutes=30,
            description="Prep 2",
            start_dt=datetime(2026, 2, 1, 8, 0),
            end_dt=datetime(2026, 2, 1, 8, 30),
            shift_assignment_id=None,
            created_by="me",
            created_at=datetime(2026, 2, 1, 8, 30),
            updated_at=None,
            updated_by=None,
        ),
    ]

    summary = _build_delay_summary(delays)

    assert summary.iloc[0]["DELAY TYPE"] == "Preparation"
    assert summary.iloc[0]["COUNT"] == 2
    assert summary.iloc[0]["TOTAL DELAY (HRS)"] == 1.5


def test_format_delay_summary_for_table_uses_percent_strings():
    summary = pd.DataFrame(
        [
            {
                "DELAY TYPE": "Equipment",
                "COUNT": 5,
                "FREQUENCY OF OCCURRENCE": 0.25,
                "TOTAL DELAY (HRS)": 27.75,
                "PERCENTAGE OF TOTAL DELAY": 0.30,
            }
        ]
    )

    formatted = _format_delay_summary_for_table(summary)

    assert formatted.iloc[0]["COUNT"] == "5"
    assert formatted.iloc[0]["FREQUENCY OF OCCURRENCE"] == "25.0%"
    assert formatted.iloc[0]["TOTAL DELAY (HRS)"] == "27.8"
    assert formatted.iloc[0]["PERCENTAGE OF TOTAL DELAY"] == "30.0%"


def test_build_excel_report_buffer_returns_workbook_bytes():
    project = _build_project()

    buffer = build_excel_report_buffer(project)

    assert buffer.getbuffer().nbytes > 0


def test_last_table_row_uses_positive_indexing():
    class RowCollection:
        def __init__(self):
            self._rows = ["header", "body"]

        def __len__(self):
            return len(self._rows)

        def __getitem__(self, index):
            if index < 0:
                raise IndexError(f"row index [{index}] out of range")
            return self._rows[index]

    class FakeTable:
        def __init__(self):
            self.rows = RowCollection()

    table = FakeTable()

    assert MillRelinePowerPointReport._last_table_row(table) == "body"


def test_clone_slide_supports_python_pptx_relationship_api():
    prs = Presentation(str(Path(__file__).resolve().parents[1] / "src" / "assets" / "mill_reline_template.pptx"))
    source_slide = prs.slides[14]

    cloned = MillRelinePowerPointReport._clone_slide(prs, source_slide)

    assert cloned is not None
    assert len(prs.slides) == 19

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "cloned.pptx"
        prs.save(output_path)
        reopened = Presentation(output_path)
        assert len(reopened.slides) == 19


def test_remap_relationship_ids_updates_embeds():
    xml = etree.fromstring(
        b"""
        <p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
          <p:blipFill>
            <a:blip xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" r:embed="rId5" />
          </p:blipFill>
        </p:pic>
        """
    )

    MillRelinePowerPointReport._remap_relationship_ids(xml, {"rId5": "rId9"})

    rel_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
    assert xml[0][0].attrib[rel_attr] == "rId9"


def test_render_template_without_figures_can_be_saved_and_reopened():
    project = _build_project()
    context = MillRelinePowerPointReport.build_context(project)
    report = MillRelinePowerPointReport()

    with tempfile.TemporaryDirectory() as tmpdir:
        output_path = Path(tmpdir) / "rendered.pptx"
        report.render(context, output_path=output_path)
        reopened = Presentation(output_path)
        assert len(reopened.slides) >= 18
