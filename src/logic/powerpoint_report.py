from __future__ import annotations

from collections.abc import Callable, Mapping, Sequence
from copy import deepcopy
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from pathlib import Path
import re
from typing import Any, BinaryIO, Optional

import pandas as pd

from logic.post_mortem import PostMortemAnalyzer
from models.phase import Phase
from models.project import Project


FigureSource = str | Path | bytes | BytesIO | BinaryIO


def _coerce_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float):
        text = f"{value:.1f}"
    else:
        text = str(value)
    text = re.sub(r"_x[0-9A-Fa-f]{4}_", "", text)
    text = text.replace("\x0b", "")
    return text


def apply_placeholder_replacements(text: str, replacements: Mapping[str, Any]) -> str:
    """
    Replace both simple placeholders like ``<PROJECT_TITLE>`` and the split
    placeholders used in the phase section like ``< | phase_num | >``.
    """
    if not text:
        return text

    result = text

    def replace_split(match: re.Match[str]) -> str:
        token = match.group(1).strip()
        if token in replacements:
            return _coerce_text(replacements[token])
        return match.group(0)

    def replace_simple(match: re.Match[str]) -> str:
        token = match.group(1).strip()
        if token in replacements:
            return _coerce_text(replacements[token])
        return match.group(0)

    result = re.sub(r"<\s*\|\s*([A-Za-z0-9_]+)\s*\|\s*>", replace_split, result)
    result = re.sub(r"<([A-Za-z0-9_]+)>", replace_simple, result)

    for key, value in replacements.items():
        result = result.replace(str(key), _coerce_text(value))

    return result


def build_default_replacements(
    project: Project,
    *,
    created_at: Optional[datetime] = None,
    extra: Optional[Mapping[str, Any]] = None,
) -> dict[str, Any]:
    created_at = created_at or datetime.now()
    project_end = project.actual_end or project.end_date or created_at
    project_start = project.actual_start or project.start_date or created_at

    total_delay_hours = None
    if project.actual_end and project.end_date:
        total_delay_hours = (project.actual_end - project.end_date).total_seconds() / 3600

    replacements: dict[str, Any] = {
        "CREATION_DATE": created_at.strftime("%B %d, %Y"),
        "PROJECT_TITLE": project.name,
        "MONTH": project_end.strftime("%B"),
        "YEAR": project_end.strftime("%Y"),
        "START_DAY": project_start.strftime("%d"),
        "END_DAY": project_end.strftime("%d"),
        "TOTAL_DELAY_HOURS": round(total_delay_hours, 1) if total_delay_hours is not None else "",
        "MILL_NAME": project.name,
    }
    if extra:
        replacements.update(extra)
    return replacements


@dataclass(frozen=True)
class FigurePlacement:
    left: float
    top: float
    width: float
    height: float


@dataclass(frozen=True)
class FigureRequest:
    name: str
    producer: Callable[..., FigureSource]
    placement: FigurePlacement
    args: tuple[Any, ...] = ()
    kwargs: Mapping[str, Any] = field(default_factory=dict)

    def render(self) -> FigureSource:
        return self.producer(*self.args, **dict(self.kwargs))


@dataclass(frozen=True)
class TableRequest:
    dataframe: pd.DataFrame
    max_rows: Optional[int] = None
    fallback_font_size_pt: Optional[float] = None
    fallback_bold: Optional[bool] = None


@dataclass(frozen=True)
class PhaseSectionContent:
    phase: Phase
    timeline_figure: Optional[FigureRequest] = None
    delayed_tasks: Optional[pd.DataFrame] = None
    replacements: Mapping[str, Any] = field(default_factory=dict)


@dataclass
class MillRelineReportContext:
    project: Project
    replacements: Mapping[str, Any] = field(default_factory=dict)
    slide_tables: Mapping[str, TableRequest] = field(default_factory=dict)
    slide_figures: Mapping[str, Sequence[FigureRequest]] = field(default_factory=dict)
    phase_sections: Sequence[PhaseSectionContent] = field(default_factory=tuple)


DEFAULT_TEMPLATE_PATH = Path("src/assets/mill_reline_template.pptx")


class MillRelinePowerPointReport:
    """
    Template-driven PowerPoint builder for mill reline post-mortem reports.

    The engine handles text replacement, table filling, and figure insertion.
    Figure creation remains fully external: callers provide callables that return
    a file path, bytes, or a binary stream consumable by ``python-pptx``.
    """

    FIXED_SLIDE_TITLES = {
        "major_delays": "Summary of Major Delays",
        "phase_summary": "Reline Phases",
        "progress": "Reline Progress",
        "phase_delay_chart": "Reline Phases & Their Delays",
        "delay_types": "Delays by Type",
        "delay_types_raw": "Raw Data",
        "inching_performance": "Inching Performance",
        "ttfi": "Time to First Inch",
    }

    DEFAULT_FIGURE_PLACEMENTS = {
        "delay_types_left": FigurePlacement(0.6, 1.6, 5.8, 4.8),
        "delay_types_right": FigurePlacement(6.7, 1.6, 5.8, 4.8),
        "progress": FigurePlacement(0.7, 1.5, 12.0, 5.2),
        "phase_delay_chart": FigurePlacement(0.7, 1.5, 12.0, 5.2),
        "inching_performance": FigurePlacement(0.6, 1.6, 8.2, 4.8),
        "shift_inching": FigurePlacement(0.7, 1.5, 12.0, 5.2),
        "ttfi": FigurePlacement(4.1, 1.7, 8.6, 4.8),
        "phase_timeline": FigurePlacement(0.6, 1.45, 12.1, 5.4),
    }
    REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    def __init__(self, template_path: Path | str = DEFAULT_TEMPLATE_PATH):
        self.template_path = Path(template_path)

    @staticmethod
    def build_phase_sections(
        project: Project,
        *,
        phase_delay_rows: int = 5,
        timeline_figure_builders: Optional[Mapping[str, Callable[[], FigureSource] | FigureRequest]] = None,
        phase_replacements: Optional[Mapping[str, Mapping[str, Any]]] = None,
    ) -> list[PhaseSectionContent]:
        sections: list[PhaseSectionContent] = []
        figure_builders = timeline_figure_builders or {}
        replacement_map = phase_replacements or {}

        for index, pid in enumerate(project.phase_order, start=1):
            phase = project.phases[pid]
            figure_builder = figure_builders.get(pid)
            figure_request: Optional[FigureRequest] = None

            if isinstance(figure_builder, FigureRequest):
                figure_request = figure_builder
            elif callable(figure_builder):
                figure_request = FigureRequest(
                    name=f"phase-{index}-timeline",
                    producer=figure_builder,
                    placement=MillRelinePowerPointReport.DEFAULT_FIGURE_PLACEMENTS["phase_timeline"],
                )

            replacements = {
                "phase_num": index,
                "phase_name": phase.name,
            }
            replacements.update(replacement_map.get(pid, {}))

            sections.append(
                PhaseSectionContent(
                    phase=phase,
                    timeline_figure=figure_request,
                    delayed_tasks=PostMortemAnalyzer.analyze_phase_delays(phase, phase_delay_rows),
                    replacements=replacements,
                )
            )
        return sections

    @staticmethod
    def build_context(
        project: Project,
        *,
        replacements: Optional[Mapping[str, Any]] = None,
        major_delay_rows: int = 8,
        phase_delay_rows: int = 5,
        slide_tables: Optional[Mapping[str, TableRequest]] = None,
        slide_figures: Optional[Mapping[str, Sequence[FigureRequest]]] = None,
        timeline_figure_builders: Optional[Mapping[str, Callable[[], FigureSource] | FigureRequest]] = None,
        phase_replacements: Optional[Mapping[str, Mapping[str, Any]]] = None,
    ) -> MillRelineReportContext:
        merged_replacements = build_default_replacements(project, extra=replacements)
        major_delays_df = PostMortemAnalyzer.major_delays(project, major_delay_rows)[["Phase", "Task", "Delay", "Notes"]].copy()
        major_delays_df.insert(3, "Type", "")

        default_tables: dict[str, TableRequest] = {
            "major_delays": TableRequest(
                dataframe=major_delays_df[["Phase", "Task", "Delay", "Type", "Notes"]],
                max_rows=major_delay_rows,
            ),
            "phase_summary": TableRequest(
                dataframe=PostMortemAnalyzer.phase_delay_summary(project),
            ),
        }

        if slide_tables:
            default_tables.update(slide_tables)

        phase_sections = MillRelinePowerPointReport.build_phase_sections(
            project,
            phase_delay_rows=phase_delay_rows,
            timeline_figure_builders=timeline_figure_builders,
            phase_replacements=phase_replacements,
        )

        return MillRelineReportContext(
            project=project,
            replacements=merged_replacements,
            slide_tables=default_tables,
            slide_figures=slide_figures or {},
            phase_sections=phase_sections,
        )

    def render(self, context: MillRelineReportContext, output_path: Optional[Path | str] = None):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ModuleNotFoundError as exc:
            raise ModuleNotFoundError(
                "python-pptx is required to render PowerPoint reports. "
                "Install it from requirements.txt before calling render()."
            ) from exc

        prs = Presentation(str(self.template_path))

        for slide in prs.slides:
            self._replace_text_on_slide(slide, context.replacements)

        for table_key, table_request in context.slide_tables.items():
            slide = self._find_slide(prs, self.FIXED_SLIDE_TITLES.get(table_key, table_key))
            if slide is None:
                continue
            table = self._find_first_table(slide)
            if table is None:
                continue
            self._populate_table(
                table,
                table_request.dataframe,
                max_rows=table_request.max_rows,
                fallback_font_size_pt=table_request.fallback_font_size_pt,
                fallback_bold=table_request.fallback_bold,
            )

        for slide_key, figure_requests in context.slide_figures.items():
            slide = self._find_slide(prs, self.FIXED_SLIDE_TITLES.get(slide_key, slide_key))
            if slide is None:
                continue
            for figure_request in figure_requests:
                picture = figure_request.render()
                self._add_picture(slide, picture, figure_request.placement, Inches)

        self._populate_phase_sections(prs, context, Inches)

        if output_path is not None:
            prs.save(str(output_path))
        return prs

    def _populate_phase_sections(self, prs, context: MillRelineReportContext, inches_fn) -> None:
        timeline_template = self._find_slide(prs, "Timeline")
        delayed_tasks_template = self._find_slide(prs, "Delayed Tasks")

        if timeline_template is None or delayed_tasks_template is None:
            return

        appendix_slide = self._find_slide(prs, "Appendix")
        insert_index = len(prs.slides) - 1
        if appendix_slide is not None:
            insert_index = self._slide_index(prs, appendix_slide)

        if not context.phase_sections:
            self._remove_slide(prs, delayed_tasks_template)
            self._remove_slide(prs, timeline_template)
            return

        timeline_slides = [timeline_template]
        delay_slides = [delayed_tasks_template]

        for _ in range(len(context.phase_sections) - 1):
            cloned_timeline = self._clone_slide(prs, timeline_template)
            self._move_slide(prs, cloned_timeline, insert_index)
            insert_index += 1
            cloned_delay = self._clone_slide(prs, delayed_tasks_template)
            self._move_slide(prs, cloned_delay, insert_index)
            insert_index += 1
            timeline_slides.append(cloned_timeline)
            delay_slides.append(cloned_delay)

        for section, timeline_slide, delay_slide in zip(context.phase_sections, timeline_slides, delay_slides):
            merged_replacements = dict(context.replacements)
            merged_replacements.update(section.replacements)
            self._replace_text_on_slide(timeline_slide, merged_replacements)
            self._replace_text_on_slide(delay_slide, merged_replacements)

            if section.timeline_figure is not None:
                self._add_picture(
                    timeline_slide,
                    section.timeline_figure.render(),
                    section.timeline_figure.placement,
                    inches_fn,
                )

            if section.delayed_tasks is not None:
                table = self._find_first_table(delay_slide)
                if table is not None:
                    phase_df = section.delayed_tasks[["Task", "Planned Duration", "Actual Duration", "Delay", "Notes"]]
                    self._populate_table(table, phase_df)

    @staticmethod
    def _slide_index(prs, slide) -> int:
        for i, candidate in enumerate(prs.slides):
            if candidate == slide:
                return i
        raise ValueError("Slide not found in presentation.")

    @staticmethod
    def _iter_text_shapes(slide):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                yield shape
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell

    def _replace_text_on_slide(self, slide, replacements: Mapping[str, Any]) -> None:
        for text_container in self._iter_text_shapes(slide):
            if not getattr(text_container, "text_frame", None):
                continue
            for paragraph in text_container.text_frame.paragraphs:
                if self._replace_phase_title_paragraph(paragraph, replacements):
                    continue
                original = paragraph.text
                updated = apply_placeholder_replacements(original, replacements)
                if updated != original:
                    self._set_paragraph_text_preserve_format(paragraph, updated)

    @staticmethod
    def _replace_phase_title_paragraph(paragraph, replacements: Mapping[str, Any]) -> bool:
        text = paragraph.text or ""
        if "phase_num" not in text or "phase_name" not in text:
            return False
        if "phase_num" not in replacements or "phase_name" not in replacements:
            return False

        phase_num = _coerce_text(replacements.get("phase_num", ""))
        phase_name = _coerce_text(replacements.get("phase_name", ""))
        if not phase_num or not phase_name:
            return False
        runs = list(paragraph.runs)

        if "Timeline" in text and len(runs) >= 5:
            runs[0].text = "Phase "
            runs[1].text = phase_num
            runs[2].text = ": "
            runs[3].text = phase_name
            runs[4].text = " - Timeline"
            for run in runs[5:]:
                run.text = ""
            return True

        if "Delayed Tasks" in text and len(runs) >= 7:
            runs[0].text = "Phase "
            runs[1].text = phase_num
            runs[2].text = ": "
            runs[3].text = phase_name
            runs[4].text = " - "
            runs[5].text = ""
            runs[6].text = "Delayed Tasks"
            for run in runs[7:]:
                run.text = ""
            return True

        return False

    @staticmethod
    def _find_slide(prs, text_fragment: str):
        fragment = text_fragment.lower()
        for slide in prs.slides:
            slide_text = " ".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
            if fragment in slide_text.lower():
                return slide
        return None

    @staticmethod
    def _find_first_table(slide):
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return shape.table
        return None

    @staticmethod
    def _last_table_row(table):
        if len(table.rows) == 0:
            raise ValueError("Cannot populate a PowerPoint table with zero rows in the template.")
        return table.rows[len(table.rows) - 1]

    @staticmethod
    def _set_paragraph_text_preserve_format(
        paragraph,
        text: str,
        *,
        fallback_font_size_pt: Optional[float] = None,
        fallback_bold: Optional[bool] = None,
    ) -> None:
        runs = list(paragraph.runs)
        if runs:
            runs[0].text = text
            for run in runs[1:]:
                run.text = ""
            return

        run = paragraph.add_run()
        run.text = text
        if fallback_font_size_pt is not None:
            from pptx.util import Pt

            run.font.size = Pt(fallback_font_size_pt)
        if fallback_bold is not None:
            run.font.bold = fallback_bold

    @classmethod
    def _set_text_frame_text_preserve_format(
        cls,
        text_frame,
        text: str,
        *,
        fallback_font_size_pt: Optional[float] = None,
        fallback_bold: Optional[bool] = None,
    ) -> None:
        paragraphs = list(text_frame.paragraphs)
        if not paragraphs:
            paragraph = text_frame.add_paragraph()
            cls._set_paragraph_text_preserve_format(
                paragraph,
                text,
                fallback_font_size_pt=fallback_font_size_pt,
                fallback_bold=fallback_bold,
            )
            return

        cls._set_paragraph_text_preserve_format(
            paragraphs[0],
            text,
            fallback_font_size_pt=fallback_font_size_pt,
            fallback_bold=fallback_bold,
        )
        for paragraph in paragraphs[1:]:
            cls._set_paragraph_text_preserve_format(
                paragraph,
                "",
                fallback_font_size_pt=fallback_font_size_pt,
                fallback_bold=fallback_bold,
            )

    @staticmethod
    def _ensure_table_row_count(table, total_rows: int) -> None:
        while len(table.rows) < total_rows:
            template_row = deepcopy(MillRelinePowerPointReport._last_table_row(table)._tr)
            table._tbl.append(template_row)

        while len(table.rows) > total_rows:
            table._tbl.remove(MillRelinePowerPointReport._last_table_row(table)._tr)

    def _populate_table(
        self,
        table,
        df: pd.DataFrame,
        max_rows: Optional[int] = None,
        fallback_font_size_pt: Optional[float] = None,
        fallback_bold: Optional[bool] = None,
    ) -> None:
        output_df = df.copy()
        if max_rows is not None:
            output_df = output_df.head(max_rows)

        total_rows = max(len(output_df) + 1, 2)
        self._ensure_table_row_count(table, total_rows)

        for row_idx in range(1, len(table.rows)):
            for col_idx in range(len(table.columns)):
                self._set_text_frame_text_preserve_format(
                    table.cell(row_idx, col_idx).text_frame,
                    "",
                    fallback_font_size_pt=fallback_font_size_pt,
                    fallback_bold=fallback_bold,
                )

        for row_offset, (_, row) in enumerate(output_df.iterrows(), start=1):
            for col_idx, value in enumerate(row.tolist()):
                if col_idx >= len(table.columns):
                    break
                self._set_text_frame_text_preserve_format(
                    table.cell(row_offset, col_idx).text_frame,
                    _coerce_text(value),
                    fallback_font_size_pt=fallback_font_size_pt,
                    fallback_bold=fallback_bold,
                )

    @staticmethod
    def _add_picture(slide, picture: FigureSource, placement: FigurePlacement, inches_fn) -> None:
        if hasattr(picture, "seek"):
            picture.seek(0)

        image_source = picture
        if isinstance(picture, bytes):
            image_source = BytesIO(picture)

        slide.shapes.add_picture(
            image_source,
            inches_fn(placement.left),
            inches_fn(placement.top),
            width=inches_fn(placement.width),
            height=inches_fn(placement.height),
        )

    @staticmethod
    def _clone_slide(prs, source_slide):
        blank_layout = prs.slide_layouts[6]
        cloned_slide = prs.slides.add_slide(blank_layout)

        rid_map: dict[str, str] = {}
        for rel in source_slide.part.rels.values():
            if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
                continue
            target = rel.target_ref if rel.is_external else rel.target_part
            new_rid = cloned_slide.part.rels._add_relationship(
                rel.reltype,
                target,
                rel.is_external,
            )
            rid_map[rel.rId] = new_rid

        for shape in source_slide.shapes:
            shape_element = deepcopy(shape.element)
            MillRelinePowerPointReport._remap_relationship_ids(shape_element, rid_map)
            cloned_slide.shapes._spTree.insert_element_before(
                shape_element, "p:extLst"
            )

        return cloned_slide

    @classmethod
    def _remap_relationship_ids(cls, element, rid_map: Mapping[str, str]) -> None:
        rel_ns = f"{{{cls.REL_NS}}}"
        for node in element.iter():
            for attr_name, attr_value in list(node.attrib.items()):
                if not attr_name.startswith(rel_ns):
                    continue
                if attr_value in rid_map:
                    node.set(attr_name, rid_map[attr_value])

    def _move_slide(self, prs, slide, new_index: int) -> None:
        current_index = self._slide_index(prs, slide)
        sld_id_lst = prs.slides._sldIdLst
        slide_id = list(sld_id_lst)[current_index]
        sld_id_lst.remove(slide_id)
        sld_id_lst.insert(new_index, slide_id)

    def _remove_slide(self, prs, slide) -> None:
        index = self._slide_index(prs, slide)
        slide_id = prs.slides._sldIdLst[index]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[index]
